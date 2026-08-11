"""Extracción de bloques de texto desde los documentos oficiales del curso.

Esta capa hace **solo una cosa**: convertir un archivo (PDF o PPTX) en una lista
plana de bloques con su número de página. No decide qué es un fragmento
recuperable —de eso se encarga `chunker.py`— ni toca la base de datos.

Están separadas a propósito: la agrupación en fragmentos es la decisión de
diseño que más impacta la calidad del RAG (Hashiyada et al. advierten que un
chunking descuidado hace que el sistema ignore el material institucional), así
que conviene poder ajustarla y testearla sin volver a parsear archivos.

**Por qué se conserva el número de página.** La tabla 17.7 exige
`pagina_inicio`/`pagina_fin`, y el cap. 8.1 pide que la exactitud factual sea
verificable «considerando el documento fuente, la página o sección
correspondiente». Sin la página, la trazabilidad se queda en el documento
completo y deja de ser auditable en la práctica.
"""

import hashlib
import re
import unicodedata
from contextlib import suppress
from dataclasses import dataclass
from pathlib import Path

# Un encabezado en PDF se detecta por tamaño de fuente relativo al cuerpo. Un
# título real suele ser sensiblemente más grande; 1.15 deja pasar las negritas
# de énfasis (que rondan el mismo tamaño) sin confundirlas con secciones.
FACTOR_TITULO = 1.15

# Un encabezado además es corto. Un párrafo en negrita grande puede superar el
# factor de tamaño, pero no suele caber en una línea de título.
MAX_PALABRAS_TITULO = 15

FORMATOS_SOPORTADOS = (".pdf", ".pptx")


class FormatoNoSoportado(ValueError):
    """El archivo no es PDF ni PPTX."""


@dataclass(frozen=True, slots=True)
class BloqueTexto:
    """Una unidad de texto tal como venía en el documento, con su procedencia.

    `es_titulo` no es cosmético: el chunker lo usa como frontera dura, porque un
    encabezado marca el inicio de un tema nuevo y unir dos secciones en un mismo
    fragmento es exactamente lo que rompe la correspondencia
    «un fragmento ↔ un objetivo de aprendizaje» del cap. 12.
    """

    texto: str
    pagina: int
    es_titulo: bool = False
    tipo: str = "texto"


def normalizar(texto: str) -> str:
    """Limpia el texto crudo del extractor sin alterar su contenido.

    Tres arreglos, todos habituales al sacar texto de un PDF:

    1. Guion de corte de línea (``proce-\\nso`` → ``proceso``). Si no se
       deshace, el término queda partido y el full-text search en español
       nunca lo encuentra, que es justamente el filtro secundario del retriever.
    2. Saltos de línea internos → espacio, porque el corte de línea del PDF es
       una decisión de maquetado, no de contenido.
    3. Espacios repetidos colapsados.

    Se normaliza además a NFC para que las tildes queden en forma compuesta:
    algunos PDF entregan ``o`` + acento combinante, que compara distinto que
    ``ó`` aunque se vea igual.
    """
    texto = unicodedata.normalize("NFC", texto)
    texto = re.sub(r"(\w)-\s*\n\s*(\w)", r"\1\2", texto)
    texto = re.sub(r"\s*\n\s*", " ", texto)
    return re.sub(r"[ \t ]+", " ", texto).strip()


def hash_archivo(ruta: Path) -> str:
    """SHA-256 del archivo, para el `hash_archivo` único de la tabla 17.6.

    Evita el error más fácil de la curación: cargar dos veces el mismo apunte
    con distinto nombre y terminar con fragmentos duplicados compitiendo en el
    retriever. Se lee por bloques porque un PDF de clase puede pesar decenas de
    megabytes.
    """
    sha = hashlib.sha256()
    with ruta.open("rb") as fh:
        for bloque in iter(lambda: fh.read(65_536), b""):
            sha.update(bloque)
    return sha.hexdigest()


def _es_titulo(texto: str, tamano: float, tamano_cuerpo: float) -> bool:
    """Heurística de encabezado: más grande que el cuerpo y corto."""
    if not texto:
        return False
    if len(texto.split()) > MAX_PALABRAS_TITULO:
        return False
    return tamano >= tamano_cuerpo * FACTOR_TITULO


def extraer_pdf(ruta: Path) -> list[BloqueTexto]:
    """PDF → bloques, detectando encabezados por tamaño de fuente.

    El tamaño del cuerpo se calcula como la **moda** de los tamaños de todo el
    documento, no como un valor fijo: cada apunte usa su propia tipografía, y
    comparar contra una constante marcaría todo como título en un documento de
    letra grande y nada en uno de letra chica.
    """
    import pymupdf

    bloques: list[BloqueTexto] = []
    with pymupdf.open(ruta) as doc:
        # Primera pasada: repartir el "peso" de cada tamaño de fuente por
        # cantidad de caracteres. Ponderar por caracteres y no por cantidad de
        # spans evita que muchos títulos cortos desplacen al cuerpo real.
        peso_por_tamano: dict[float, int] = {}
        for pagina in doc:
            for bloque in pagina.get_text("dict")["blocks"]:
                for linea in bloque.get("lines", []):
                    for span in linea.get("spans", []):
                        tam = round(span["size"], 1)
                        peso_por_tamano[tam] = peso_por_tamano.get(tam, 0) + len(span["text"])

        if not peso_por_tamano:
            return []
        tamano_cuerpo = max(peso_por_tamano, key=lambda t: peso_por_tamano[t])

        # Segunda pasada: armar los bloques ya clasificados.
        for num_pagina, pagina in enumerate(doc, start=1):
            for bloque in pagina.get_text("dict")["blocks"]:
                lineas = bloque.get("lines", [])
                if not lineas:
                    continue  # bloque de imagen: sin texto que indexar

                spans = [s for ln in lineas for s in ln.get("spans", [])]
                if not spans:
                    continue

                # Los spans de una misma línea se unen sin separador (un span
                # nuevo solo marca un cambio de tipografía, que puede caer a
                # mitad de palabra), pero las líneas se unen con "\n" y no con
                # "" — si no, la última palabra de un renglón queda pegada a la
                # primera del siguiente ("para" + "reducir" → "parareducir").
                # `normalizar` convierte después ese "\n" en espacio y de paso
                # deshace los guiones de corte.
                texto = normalizar(
                    "\n".join(
                        "".join(s["text"] for s in ln.get("spans", []))
                        for ln in lineas
                        if ln.get("spans")
                    )
                )
                if not texto:
                    continue

                tamano = max(s["size"] for s in spans)
                bloques.append(
                    BloqueTexto(
                        texto=texto,
                        pagina=num_pagina,
                        es_titulo=_es_titulo(texto, tamano, tamano_cuerpo),
                    )
                )

    return bloques


def _tabla_a_texto(tabla) -> str:
    """Aplana una tabla de PowerPoint a texto con filas separadas por '|'.

    El retriever es textual: una tabla que no se serializa a texto es
    invisible para el full-text search y para el LLM. Se conserva la estructura
    con separadores en vez de perderla, porque para un perfil visual la tabla
    comparativa es justamente el recurso que pide la tabla 11.1.
    """
    filas = []
    for fila in tabla.rows:
        celdas = [normalizar(c.text) for c in fila.cells]
        if any(celdas):
            filas.append(" | ".join(celdas))
    return "\n".join(filas)


def extraer_pptx(ruta: Path) -> list[BloqueTexto]:
    """PPTX → bloques, una diapositiva por "página".

    En una presentación la diapositiva ya es la unidad temática que en un PDF
    hay que inferir: tiene título propio y un cuerpo acotado. Por eso el título
    del placeholder se marca como encabezado directamente, sin heurística de
    tamaño.
    """
    from pptx import Presentation

    bloques: list[BloqueTexto] = []
    presentacion = Presentation(str(ruta))

    for num_diapo, diapo in enumerate(presentacion.slides, start=1):
        forma_titulo = None
        # Hay layouts sin placeholder de título; no es un error.
        with suppress(AttributeError, ValueError):
            forma_titulo = diapo.shapes.title

        if forma_titulo is not None and normalizar(forma_titulo.text):
            bloques.append(
                BloqueTexto(
                    texto=normalizar(forma_titulo.text),
                    pagina=num_diapo,
                    es_titulo=True,
                )
            )

        for forma in diapo.shapes:
            if forma is forma_titulo:
                continue

            if getattr(forma, "has_table", False):
                texto = _tabla_a_texto(forma.table)
                if texto:
                    bloques.append(
                        BloqueTexto(texto=texto, pagina=num_diapo, tipo="tabla")
                    )
                continue

            if not getattr(forma, "has_text_frame", False):
                continue

            # Cada párrafo por separado: en una diapositiva de viñetas, cada
            # viñeta es una idea independiente y unirlas produce un bloque
            # corrido que el chunker ya no puede separar por tema.
            for parrafo in forma.text_frame.paragraphs:
                texto = normalizar(parrafo.text)
                if texto:
                    bloques.append(BloqueTexto(texto=texto, pagina=num_diapo))

    return bloques


def extraer(ruta: Path) -> list[BloqueTexto]:
    """Despacha al extractor según la extensión del archivo."""
    sufijo = ruta.suffix.lower()
    if sufijo == ".pdf":
        return extraer_pdf(ruta)
    if sufijo == ".pptx":
        return extraer_pptx(ruta)
    raise FormatoNoSoportado(
        f"'{ruta.name}': formato {sufijo or '(sin extensión)'} no soportado. "
        f"Formatos válidos: {', '.join(FORMATOS_SOPORTADOS)}"
    )
